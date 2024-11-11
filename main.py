# Part 1: UI and Setup (save as photo_organizer_ui.py)

import os
import re
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from PIL import Image
import pytesseract
from pptx import Presentation
from pptx.util import Inches
import threading
from collections import defaultdict
import datetime

class PhotoOrganizerApp:
    def __init__(self, root):
        self.root = root
        self.root.title("Photo Collection Organizer")
        self.root.geometry("1100x700")
        
        # Configure colors for elegant dark theme
        self.colors = {
            'primary': '#BB86FC',      # Elegant Purple
            'secondary': '#121212',    # Very Dark Gray
            'background': '#1E1E1E',   # Dark Gray
            'surface': '#242424',      # Slightly Lighter Gray
            'text': '#FFFFFF',         # White
            'text_secondary': '#B3B3B3', # Light Gray
            'accent': '#03DAC6',       # Teal Accent
            'error': '#CF6679',        # Error Red
            'success': '#4CAF50',      # Elegant Green
            'entry_text': '#00E676'    # Bright Elegant Green
        }
        
        # Configure root background
        self.root.configure(bg=self.colors['background'])
        
        # Configure styles for dark theme
        self.style = ttk.Style()
        self.style.configure('TFrame', background=self.colors['background'])
        self.style.configure('Surface.TFrame', background=self.colors['surface'])
        
        # Configure Label styles
        self.style.configure('TLabel', 
                           background=self.colors['background'],
                           foreground=self.colors['text'])
        
        self.style.configure('Surface.TLabel', 
                           background=self.colors['surface'],
                           foreground=self.colors['text'])
        
        self.style.configure('Header.TLabel', 
                           background=self.colors['background'],
                           foreground=self.colors['primary'],
                           font=('Helvetica', 28, 'bold'))
        
        self.style.configure('SubHeader.TLabel', 
                           background=self.colors['background'],
                           foreground=self.colors['text_secondary'],
                           font=('Helvetica', 14))
        
        # Configure LabelFrame style
        self.style.configure('TLabelframe', 
                           background=self.colors['surface'],
                           foreground=self.colors['text'])
        
        self.style.configure('TLabelframe.Label', 
                           background=self.colors['surface'],
                           foreground=self.colors['primary'],
                           font=('Helvetica', 11, 'bold'))
        
        # Configure Entry style
        self.style.configure('TEntry', 
                           fieldbackground=self.colors['surface'],
                           foreground=self.colors['entry_text'],
                           insertcolor=self.colors['entry_text'],
                           font=('Helvetica', 11, 'bold'))
        
        # Configure Button styles
        self.style.configure('TButton', 
                           background=self.colors['primary'],
                           font=('Helvetica', 10))
        
        self.style.configure('Process.TButton',
                           background=self.colors['primary'],
                           font=('Helvetica', 14, 'bold'),
                           padding=15)
        
        # Configure Progressbar style
        self.style.configure('Horizontal.TProgressbar',
                           background=self.colors['primary'],
                           troughcolor=self.colors['surface'])
        
        # Statistics Label styles
        self.style.configure('Stats.TLabel',
                           background=self.colors['surface'],
                           foreground=self.colors['accent'],
                           font=('Helvetica', 24, 'bold'))
        
        self.style.configure('StatsText.TLabel',
                           background=self.colors['surface'],
                           foreground=self.colors['text_secondary'],
                           font=('Helvetica', 12))
        
        # Main container
        self.main_frame = ttk.Frame(root, style='Main.TFrame')
        self.main_frame.pack(fill=tk.BOTH, expand=True)
        
        # Left Panel (Input/Output Configuration)
        left_panel = ttk.Frame(self.main_frame, style='Main.TFrame')
        left_panel.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=20, pady=20)
        
        # Header
        ttk.Label(left_panel, 
                 text="Photo Collection Organizer", 
                 style='Header.TLabel').pack(pady=(0, 5))
        
        ttk.Label(left_panel,
                 text="Organize your photos into PowerPoint presentations automatically",
                 style='SubHeader.TLabel').pack(pady=(0, 20))
        
        # Input Configuration
        input_frame = ttk.LabelFrame(left_panel, 
                                   text="Input Configuration",
                                   padding="15")
        input_frame.pack(fill=tk.X, pady=(0, 15))
        
        ttk.Label(input_frame, 
                 text="Source Folder:",
                 font=('Helvetica', 10, 'bold')).pack(anchor='w')
        
        input_browse_frame = ttk.Frame(input_frame)
        input_browse_frame.pack(fill=tk.X, pady=(5, 10))
        
        self.input_path = tk.StringVar()
        self.input_entry = tk.Entry(input_browse_frame, 
                                   textvariable=self.input_path,
                                   width=45,
                                   font=('Helvetica', 11, 'bold'),
                                   bg=self.colors['surface'],
                                   fg=self.colors['entry_text'],
                                   insertbackground=self.colors['entry_text'])
        self.input_entry.pack(side=tk.LEFT, padx=(0, 10))
        
        ttk.Button(input_browse_frame,
                  text="Browse",
                  style='Custom.TButton',
                  command=self.browse_input).pack(side=tk.LEFT)
        
        # Output Configuration
        output_frame = ttk.LabelFrame(left_panel,
                                    text="Output Configuration",
                                    padding="15")
        output_frame.pack(fill=tk.X, pady=(0, 15))
        
        ttk.Label(output_frame,
                 text="Output Folder:",
                 font=('Helvetica', 10, 'bold')).pack(anchor='w')
        
        output_browse_frame = ttk.Frame(output_frame)
        output_browse_frame.pack(fill=tk.X, pady=(5, 10))
        
        self.output_path = tk.StringVar()
        self.output_entry = tk.Entry(output_browse_frame,
                                    textvariable=self.output_path,
                                    width=45,
                                    font=('Helvetica', 11, 'bold'),
                                    bg=self.colors['surface'],
                                    fg=self.colors['entry_text'],
                                    insertbackground=self.colors['entry_text'])
        self.output_entry.pack(side=tk.LEFT, padx=(0, 10))
        
        ttk.Button(output_browse_frame,
                  text="Browse",
                  style='Custom.TButton',
                  command=self.browse_output).pack(side=tk.LEFT)
        
        # Right Panel (Statistics and Progress)
        right_panel = ttk.Frame(self.main_frame, style='Surface.TFrame')
        right_panel.pack(side=tk.RIGHT, fill=tk.BOTH, padx=20, pady=20, ipadx=20, ipady=20)
        
        # Statistics
        stats_frame = ttk.Frame(right_panel, style='Surface.TFrame')
        stats_frame.pack(fill=tk.X, pady=(20, 30))
        
        # Initialize stats_widgets list
        self.stats_widgets = []
        
        # Create statistics widgets and store references
        self.create_stat_widget(stats_frame, "Total Photos", "0", 0)
        self.create_stat_widget(stats_frame, "Total Slides", "0", 1)
        self.create_stat_widget(stats_frame, "Total Groups", "0", 2)
        self.create_stat_widget(stats_frame, "Unmatched", "0", 3)
        
        # Progress Bar
        progress_frame = ttk.Frame(right_panel, style='Surface.TFrame')
        progress_frame.pack(fill=tk.X, pady=(0, 20))
        
        self.progress_var = tk.DoubleVar()
        self.progress_bar = ttk.Progressbar(progress_frame,
                                          variable=self.progress_var,
                                          maximum=100)
        self.progress_bar.pack(fill=tk.X, padx=10)
        
        self.status_label = ttk.Label(progress_frame,
                                    text="Ready to process",
                                    style='StatsText.TLabel')
        self.status_label.pack(pady=(10, 0))
        
        # Details Button
        self.details_button = ttk.Button(right_panel,
                                       text="View Processing Details",
                                       style='Custom.TButton',
                                       command=self.show_details_window,
                                       state='disabled')  # Initially disabled
        self.details_button.pack(pady=(0, 10))
        
        # Store processing details as instance variable
        self.processing_details = None
        
        # Process Button
        self.process_button = ttk.Button(right_panel,
                                       text="▶ Start Processing",
                                       style='Process.TButton',
                                       command=self.start_processing)
        self.process_button.pack(pady=20)

    def create_stat_widget(self, parent, label, value, column):
        """Create a statistics widget and store reference"""
        frame = ttk.Frame(parent, style='Surface.TFrame')
        frame.grid(row=0, column=column, padx=15)
        
        # Create and store reference to value label
        value_label = ttk.Label(frame,
                              text=value,
                              style='Stats.TLabel')
        value_label.pack()
        
        ttk.Label(frame,
                 text=label,
                 style='StatsText.TLabel').pack()
        
        # Add to stats_widgets list
        self.stats_widgets.append(value_label)
        
        return value_label

    def browse_input(self):
        folder = filedialog.askdirectory(title="Select Input Folder")
        if folder:
            self.input_path.set(folder)
            self.input_entry.delete(0, tk.END)  # Clear current text
            self.input_entry.insert(0, folder)  # Insert new path

    def browse_output(self):
        folder = filedialog.askdirectory(title="Select Output Folder")
        if folder:
            self.output_path.set(folder)
            self.output_entry.delete(0, tk.END)  # Clear current text
            self.output_entry.insert(0, folder)  # Insert new path

    def validate_inputs(self):
        if not self.input_path.get() or not self.output_path.get():
            messagebox.showerror("Error", "Please select both input and output folders")
            return False
        if not os.path.exists(self.input_path.get()):
            messagebox.showerror("Error", "Input folder does not exist")
            return False
        return True

    def start_processing(self):
        if not self.validate_inputs():
            return
        
        self.process_button.configure(state='disabled')
        self.progress_var.set(0)
        self.status_label.configure(text="Processing...")
        
        # Start processing in a separate thread
        thread = threading.Thread(target=self.process_photos)
        thread.start()

    def reset_ui(self):
        self.process_button.configure(state='normal')
        self.status_label.configure(text="Ready")
        self.progress_var.set(100)

    def update_statistics(self, total_photos, total_slides, total_groups, unmatched_photos):
        """Update the statistics display"""
        values = [total_photos, total_slides, total_groups, unmatched_photos]
        for widget, value in zip(self.stats_widgets, values):
            widget.configure(text=str(value))

    def process_photos(self):
        try:
            input_folder = self.input_path.get()
            output_folder = self.output_path.get()
            
            # Get all image files
            image_files = [f for f in os.listdir(input_folder) 
                         if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp'))]
            
            if not image_files:
                self.root.after(0, lambda: messagebox.showwarning("Warning", "No image files found"))
                return
            
            # Group photos by their OCR-extracted pattern
            photo_groups = defaultdict(list)
            # Modified patterns to allow 7 or more digits
            patterns = [
                r'([1-2])[-_]?(\d{9,})',  # Matches 7 or more digits
                
            ]
            
            # Add tracking for unmatched photos
            unmatched_photos = []
            
            # Process each image with OCR
            for idx, img_file in enumerate(image_files):
                try:
                    img_path = os.path.join(input_folder, img_file)
                    # Open and preprocess image
                    image = Image.open(img_path)
                    
                    # Convert to RGB if necessary
                    if image.mode != 'RGB':
                        image = image.convert('RGB')
                    
                    # Enhance image for better OCR
                    from PIL import ImageEnhance
                    enhancer = ImageEnhance.Contrast(image)
                    image = enhancer.enhance(2.0)  # Increase contrast
                    # Extract text using Tesseract with custom configuration
                    custom_config = r'--oem 3 --psm 6'  # Added closing quote
                    text = pytesseract.image_to_string(image, config=custom_config)
                    print(f"OCR Text for {img_file}: {text}")  # Debug print
                    
                    # Try all patterns
                    pattern_found = False
                    for pattern in patterns:
                        match = re.search(pattern, text)
                        if match:
                            key = f"{match.group(1)}-{match.group(2)}"  # Standardize format
                            photo_groups[key].append(img_file)
                            print(f"Found pattern {key} in {img_file}")  # Debug print
                            pattern_found = True
                            break
                    
                    if not pattern_found:
                        # Try finding pattern in filename as fallback
                        for pattern in patterns:
                            match = re.search(pattern, img_file)
                            if match:
                                key = f"{match.group(1)}-{match.group(2)}"
                                photo_groups[key].append(img_file)
                                print(f"Found pattern {key} in filename {img_file}")
                                pattern_found = True
                                break
                    
                    if not pattern_found:
                        unmatched_photos.append(img_file)
                        print(f"No pattern match found in {img_file}")
                    
                    # Update progress
                    progress = (idx / len(image_files)) * 50
                    self.root.after(0, lambda p=progress: self.progress_var.set(p))
                    
                except Exception as e:
                    print(f"Error processing {img_file}: {str(e)}")
                    unmatched_photos.append(img_file)
                    continue
            
            print(f"Total groups found: {len(photo_groups)}")  # Debug print
            print(f"Groups: {dict(photo_groups)}")  # Debug print

            # Create PowerPoint
            prs = Presentation()
            total_slides = 0
            
            # Process each group
            total_groups = len(photo_groups)
            for idx, (group_key, photos) in enumerate(photo_groups.items(), 1):
                # Create one slide for all photos in this group
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
                
                # Use full slide dimensions with small margins
                margin = Inches(0.2)  # 0.2 inch margin
                slide_width = Inches(10) - (2 * margin)
                slide_height = Inches(7.5) - (2 * margin)
                
                # Calculate photo width to fit all photos in one row
                num_photos = len(photos)
                photo_width = slide_width / num_photos
                photo_height = slide_height  # Full height of slide
                
                # Add photos to slide in a single row
                for photo_idx, photo in enumerate(photos):
                    # Calculate x position for each photo
                    x = margin + (photo_idx * photo_width)
                    y = margin
                    
                    try:
                        img_path = os.path.join(input_folder, photo)
                        slide.shapes.add_picture(img_path, x, y, photo_width, photo_height)
                    except Exception as e:
                        print(f"Error adding photo {photo}: {str(e)}")
                
                total_slides += 1
                
                # Update progress
                progress = 50 + (idx / total_groups) * 50
                self.root.after(0, lambda p=progress: self.progress_var.set(p))
            
            # Save presentation with error handling
            output_path = os.path.join(output_folder, "organized_photos.pptx")
            try:
                prs.save(output_path)
            except PermissionError:
                # Try saving with a different filename
                timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
                alternative_path = os.path.join(output_folder, f"organized_photos_{timestamp}.pptx")
                try:
                    prs.save(alternative_path)
                    output_path = alternative_path  # Update output_path for success message
                except Exception as save_error:
                    raise Exception(f"Could not save presentation. Please ensure PowerPoint is closed and you have write permissions. Error: {str(save_error)}")
            
            # Create summary text
            summary = f"""Processing Summary:
• Total Photos: {len(image_files)}
• Successfully Grouped Photos: {len(image_files) - len(unmatched_photos)}
• Unmatched Photos: {len(unmatched_photos)}
• Total Groups Created: {len(photo_groups)}
• Groups Details:"""

            # Add group details
            for group_key, photos in photo_groups.items():
                summary += f"\n  - Group {group_key}: {len(photos)} photos"

            # Add unmatched photos list if any
            if unmatched_photos:
                summary += "\n\nUnmatched Photos:"
                for photo in unmatched_photos:
                    summary += f"\n  • {photo}"

            # Update statistics including unmatched
            self.root.after(0, lambda: self.update_statistics(
                len(image_files),          # total photos
                len(photo_groups),         # total slides (one per group)
                len(photo_groups),         # total groups
                len(unmatched_photos)      # unmatched photos
            ))
            
            # Store the processing details
            self.processing_details = {
                'groups': dict(photo_groups),
                'unmatched': unmatched_photos,
                'stats': {
                    'total_photos': len(image_files),
                    'grouped_photos': len(image_files) - len(unmatched_photos),
                    'total_groups': len(photo_groups),
                    'unmatched_count': len(unmatched_photos)
                }
            }

            # Enable the details button after processing
            self.root.after(0, lambda: self.details_button.configure(state='normal'))
            
            # Show simple success message
            self.root.after(0, lambda: self.show_success(output_path))

        except Exception as e:
            error_message = str(e)  # Capture the error message
            self.root.after(0, lambda: messagebox.showerror("Error", f"An error occurred: {error_message}"))
        
        finally:
            self.root.after(0, self.reset_ui)

    def show_success(self, output_path):
        """Show success message in dark theme"""
        success_window = tk.Toplevel(self.root)
        success_window.title("Success")
        success_window.configure(bg=self.colors['background'])
        success_window.geometry("400x250")
        
        # Make window match dark theme
        success_window.transient(self.root)
        success_window.grab_set()
        
        # Create frame for content
        content_frame = ttk.Frame(success_window, style='Surface.TFrame')
        content_frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=20)
        
        # Success icon (checkmark)
        ttk.Label(content_frame,
                 text="✓",
                 style='Stats.TLabel',
                 font=('Helvetica', 48)).pack(pady=(10, 0))
        
        # Success message
        ttk.Label(content_frame,
                 text="Processing Complete!",
                 style='Header.TLabel',
                 font=('Helvetica', 16, 'bold')).pack(pady=(10, 5))
        
        # File path (with word wrap)
        path_label = ttk.Label(content_frame,
                             text=f"Presentation saved as:\n{output_path}",
                             style='Surface.TLabel',
                             wraplength=350,
                             justify='center')
        path_label.pack(pady=(0, 20))
        
        # OK button
        ttk.Button(content_frame,
                  text="OK",
                  style='Custom.TButton',
                  command=success_window.destroy).pack()
        
        # Center the window
        success_window.update_idletasks()
        x = self.root.winfo_x() + (self.root.winfo_width() - success_window.winfo_width()) // 2
        y = self.root.winfo_y() + (self.root.winfo_height() - success_window.winfo_height()) // 2
        success_window.geometry(f"+{x}+{y}")
        
        # Bring to front
        success_window.lift()
        success_window.focus_force()

    def show_details_window(self):
        """Show detailed information in a scrollable window"""
        if not self.processing_details:
            return

        details_window = tk.Toplevel(self.root)
        details_window.title("Processing Details")
        details_window.configure(bg=self.colors['background'])
        details_window.geometry("800x600")

        # Main frame
        main_frame = ttk.Frame(details_window, style='Surface.TFrame')
        main_frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=20)

        # Create notebook for tabbed view
        notebook = ttk.Notebook(main_frame)
        notebook.pack(fill=tk.BOTH, expand=True)

        # Summary Tab
        summary_frame = ttk.Frame(notebook, style='Surface.TFrame')
        notebook.add(summary_frame, text='Summary')

        stats = self.processing_details['stats']
        groups = self.processing_details['groups']
        
        # Create detailed summary text
        summary_text = "📊 Processing Summary\n"
        summary_text += "═══════════════════\n\n"
        summary_text += f"Total Photos: {stats['total_photos']}\n"
        summary_text += f"Successfully Grouped: {stats['grouped_photos']}\n"
        summary_text += f"Total Groups: {stats['total_groups']}\n"
        summary_text += f"Unmatched Photos: {stats['unmatched_count']}\n\n"
        
        # Add group summary
        summary_text += "📁 Groups Overview\n"
        summary_text += "═══════════════\n\n"
        for group_key, photos in groups.items():
            summary_text += f"Group {group_key}: {len(photos)} photos\n"
        
        # Add unmatched photos if any exist
        if self.processing_details['unmatched']:
            summary_text += "\n❌ Unmatched Photos\n"
            summary_text += "═══════════════════\n\n"
            for idx, photo in enumerate(self.processing_details['unmatched'], 1):
                summary_text += f"{idx}. {photo}\n"
        
        self.create_scrolled_text(summary_frame, summary_text)

        # Detailed Groups Tab
        groups_frame = ttk.Frame(notebook, style='Surface.TFrame')
        notebook.add(groups_frame, text='Group Details')

        groups_text = "📑 Detailed Group Contents\n"
        groups_text += "═════════════════════════\n\n"
        
        # Sort groups by key for better organization
        for group_key in sorted(groups.keys()):
            photos = groups[group_key]
            groups_text += f"Group {group_key} ({len(photos)} photos):\n"
            groups_text += "─" * 40 + "\n"  # Separator line
            for idx, photo in enumerate(photos, 1):
                groups_text += f"  {idx}. {photo}\n"
            groups_text += "\n"
            
        self.create_scrolled_text(groups_frame, groups_text)

        # Center the window
        details_window.update_idletasks()
        x = self.root.winfo_x() + (self.root.winfo_width() - details_window.winfo_width()) // 2
        y = self.root.winfo_y() + (self.root.winfo_height() - details_window.winfo_height()) // 2
        details_window.geometry(f"+{x}+{y}")

    def create_scrolled_text(self, parent, content):
        """Helper method to create scrolled text widget"""
        frame = ttk.Frame(parent, style='Surface.TFrame')
        frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)

        # Create text widget with scrollbar
        text_widget = tk.Text(frame,
                            wrap=tk.WORD,
                            font=('Consolas', 11),
                            bg=self.colors['surface'],
                            fg=self.colors['text'],
                            relief=tk.FLAT)
        scrollbar = ttk.Scrollbar(frame, orient="vertical", command=text_widget.yview)
        text_widget.configure(yscrollcommand=scrollbar.set)

        # Pack widgets
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        text_widget.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        # Insert content and disable editing
        text_widget.insert('1.0', content)
        text_widget.configure(state='disabled')

        return text_widget

if __name__ == "__main__":
    # Set Tesseract path if it's not in system PATH
    # pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'  # Uncomment and modify for Windows
    
    root = tk.Tk()
    app = PhotoOrganizerApp(root)
    root.mainloop()        